"""
Enhanced Integration Module for the Proposal Builder System
"""
import asyncio
import uuid
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Union
from enum import Enum
from dataclasses import dataclass, asdict
import json
import base64
from pathlib import Path
import aiofiles
import httpx
from pydantic import BaseModel, Field, validator
import jinja2
from weasyprint import HTML, CSS
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from io import BytesIO


class ProposalType(str, Enum):
    """Proposal document types"""
    EXECUTIVE_SUMMARY = "executive_summary"
    TECHNICAL_PROPOSAL = "technical_proposal"
    COMPETITIVE_ANALYSIS = "competitive_analysis"
    POC_PLAYBOOK = "poc_playbook"
    COMPLIANCE_ANNEX = "compliance_annex"
    COST_ANALYSIS = "cost_analysis"


class ExportFormat(str, Enum):
    """Export formats"""
    HTML = "html"
    PDF = "pdf"
    PPTX = "pptx"
    DOCX = "docx"
    JSON = "json"


class ProposalStatus(str, Enum):
    """Proposal status"""
    DRAFT = "draft"
    IN_REVIEW = "in_review"
    APPROVED = "approved"
    REJECTED = "rejected"
    PUBLISHED = "published"


@dataclass
class ClientContext:
    """Client context for proposal generation"""
    name: str
    sector: str
    size: str  # SME, Enterprise, Government
    region: str
    current_stack: List[str]
    pain_points: List[str]
    compliance_requirements: List[str]
    budget_range: str
    timeline: str
    decision_makers: List[Dict[str, str]]
    competitors: List[str]
    kpis: List[Dict[str, Any]]


@dataclass
class TechnicalRequirements:
    """Technical requirements for proposal"""
    scenario: str
    primary_stack: str  # IBM, Azure, Google, AWS, NVIDIA
    secondary_stacks: List[str]
    deployment_model: str  # cloud, hybrid, on-premise
    integration_points: List[str]
    security_requirements: List[str]
    performance_requirements: Dict[str, Any]
    scalability_requirements: Dict[str, Any]


@dataclass
class CompetitiveAnalysis:
    """Competitive analysis data"""
    dimensions: List[str]
    competitors: List[str]
    scores: Dict[str, Dict[str, float]]  # competitor -> dimension -> score
    weights: Dict[str, float]  # dimension -> weight
    rationale: Dict[str, Dict[str, str]]  # competitor -> dimension -> explanation


class ProposalTemplate:
    """Template management for proposals"""
    
    def __init__(self, template_dir: str = "templates"):
        self.template_dir = Path(template_dir)
        self.jinja_env = jinja2.Environment(
            loader=jinja2.FileSystemLoader(template_dir),
            autoescape=jinja2.select_autoescape(['html', 'xml'])
        )
        
        # Arabic/English support
        self.setup_i18n()
    
    def setup_i18n(self):
        """Setup internationalization"""
        self.translations = {
            'ar': {
                'executive_summary': '?????? ????????',
                'technical_proposal': '??????? ??????',
                'competitive_analysis': '??????? ????????',
                'compliance_requirements': '??????? ????????',
                'cost_benefit_analysis': '????? ??????? ????????',
                'implementation_timeline': '?????? ?????? ???????',
                'risk_assessment': '????? ???????',
                'recommendations': '????????'
            },
            'en': {
                'executive_summary': 'Executive Summary',
                'technical_proposal': 'Technical Proposal', 
                'competitive_analysis': 'Competitive Analysis',
                'compliance_requirements': 'Compliance Requirements',
                'cost_benefit_analysis': 'Cost-Benefit Analysis',
                'implementation_timeline': 'Implementation Timeline',
                'risk_assessment': 'Risk Assessment',
                'recommendations': 'Recommendations'
            }
        }
    
    def get_template(self, template_name: str) -> jinja2.Template:
        """Get Jinja2 template"""
        return self.jinja_env.get_template(f"{template_name}.html")
    
    def render_template(
        self, 
        template_name: str, 
        context: Dict[str, Any],
        language: str = 'en'
    ) -> str:
        """Render template with context"""
        template = self.get_template(template_name)
        
        # Add translations to context
        context['t'] = self.translations.get(language, self.translations['en'])
        context['language'] = language
        context['is_rtl'] = language == 'ar'
        
        return template.render(**context)


class ProposalGenerator:
    """Enhanced proposal generation with AI integration"""
    
    def __init__(self, template_manager: ProposalTemplate):
        self.template_manager = template_manager
        self.ai_client = None  # Would integrate with local LLM
        
    async def generate_executive_summary(
        self, 
        client_context: ClientContext,
        tech_requirements: TechnicalRequirements,
        language: str = 'en'
    ) -> str:
        """Generate executive summary"""
        
        # AI-enhanced content generation (placeholder)
        context = {
            'client': asdict(client_context),
            'technology': asdict(tech_requirements),
            'generated_at': datetime.now().isoformat(),
            'key_benefits': self._generate_key_benefits(client_context, tech_requirements),
            'roi_projection': self._calculate_roi_projection(client_context),
            'implementation_phases': self._generate_implementation_phases(tech_requirements)
        }
        
        return self.template_manager.render_template(
            'executive_summary', 
            context, 
            language
        )
    
    async def generate_technical_proposal(
        self,
        client_context: ClientContext,
        tech_requirements: TechnicalRequirements,
        language: str = 'en'
    ) -> str:
        """Generate detailed technical proposal"""
        
        context = {
            'client': asdict(client_context),
            'technology': asdict(tech_requirements),
            'architecture_diagram': await self._generate_architecture_diagram(tech_requirements),
            'integration_matrix': self._create_integration_matrix(tech_requirements),
            'security_framework': self._generate_security_framework(tech_requirements),
            'performance_benchmarks': await self._generate_performance_benchmarks(tech_requirements),
            'scalability_plan': self._generate_scalability_plan(tech_requirements)
        }
        
        return self.template_manager.render_template(
            'technical_proposal',
            context,
            language
        )
    
    async def generate_competitive_analysis(
        self,
        competitive_data: CompetitiveAnalysis,
        client_context: ClientContext,
        language: str = 'en'
    ) -> str:
        """Generate competitive analysis with scoring"""
        
        # Calculate weighted scores
        final_scores = self._calculate_competitive_scores(competitive_data)
        
        # Generate comparison charts
        charts = await self._generate_competitive_charts(competitive_data)
        
        context = {
            'competitive_data': asdict(competitive_data),
            'final_scores': final_scores,
            'charts': charts,
            'recommendations': self._generate_competitive_recommendations(
                competitive_data, 
                final_scores
            ),
            'risk_mitigation': self._generate_risk_mitigation_strategies(competitive_data)
        }
        
        return self.template_manager.render_template(
            'competitive_analysis',
            context,
            language
        )
    
    async def generate_compliance_annex(
        self,
        client_context: ClientContext,
        compliance_mapping: Dict[str, Any],
        language: str = 'en'
    ) -> str:
        """Generate compliance annex with regulatory mapping"""
        
        context = {
            'client': asdict(client_context),
            'compliance_mapping': compliance_mapping,
            'regulatory_frameworks': self._get_applicable_frameworks(client_context),
            'control_matrix': await self._generate_control_matrix(compliance_mapping),
            'gap_analysis': self._perform_gap_analysis(compliance_mapping),
            'remediation_plan': self._generate_remediation_plan(compliance_mapping)
        }
        
        return self.template_manager.render_template(
            'compliance_annex',
            context,
            language
        )
    
    def _generate_key_benefits(
        self, 
        client_context: ClientContext,
        tech_requirements: TechnicalRequirements
    ) -> List[Dict[str, str]]:
        """Generate key benefits based on client context"""
        
        benefits = []
        
        # Sector-specific benefits
        if client_context.sector.lower() == 'banking':
            benefits.extend([
                {
                    'title': 'Enhanced Regulatory Compliance',
                    'description': 'Automated SAMA compliance monitoring and reporting',
                    'impact': 'Reduces compliance costs by 40%'
                },
                {
                    'title': 'Improved Risk Management',
                    'description': 'Real-time risk assessment and mitigation',
                    'impact': 'Reduces operational risk by 60%'
                }
            ])
        elif client_context.sector.lower() == 'government':
            benefits.extend([
                {
                    'title': 'Citizen Service Excellence',
                    'description': 'Streamlined digital services with AI assistance',
                    'impact': 'Improves citizen satisfaction by 45%'
                },
                {
                    'title': 'Operational Efficiency',
                    'description': 'Automated processes and intelligent document processing',
                    'impact': 'Reduces processing time by 70%'
                }
            ])
        
        # Technology-specific benefits
        if tech_requirements.primary_stack.lower() == 'ibm':
            benefits.append({
                'title': 'Enterprise-Grade AI',
                'description': 'IBM Watson AI with industry-specific models',
                'impact': 'Increases decision accuracy by 85%'
            })
        
        return benefits
    
    def _calculate_roi_projection(self, client_context: ClientContext) -> Dict[str, Any]:
        """Calculate ROI projection"""
        
        # Simplified ROI calculation
        base_savings = 1000000  # Base annual savings
        
        # Adjust based on company size
        size_multiplier = {
            'SME': 0.3,
            'Enterprise': 1.0,
            'Government': 1.5
        }.get(client_context.size, 1.0)
        
        annual_savings = base_savings * size_multiplier
        implementation_cost = annual_savings * 0.8  # 80% of annual savings
        
        return {
            'implementation_cost': implementation_cost,
            'annual_savings': annual_savings,
            'payback_period_months': 12,
            'three_year_roi': ((annual_savings * 3 - implementation_cost) / implementation_cost) * 100,
            'break_even_month': 10
        }
    
    def _generate_implementation_phases(
        self, 
        tech_requirements: TechnicalRequirements
    ) -> List[Dict[str, Any]]:
        """Generate implementation phases"""
        
        phases = [
            {
                'name': 'Assessment & Planning',
                'duration_weeks': 2,
                'key_activities': [
                    'Current state assessment',
                    'Requirements validation',
                    'Architecture design',
                    'Risk assessment'
                ],
                'deliverables': [
                    'Assessment report',
                    'Implementation plan',
                    'Risk mitigation strategy'
                ]
            },
            {
                'name': 'Foundation Setup',
                'duration_weeks': 4,
                'key_activities': [
                    'Infrastructure provisioning',
                    'Security configuration',
                    'Base platform installation',
                    'Integration framework setup'
                ],
                'deliverables': [
                    'Configured infrastructure',
                    'Security documentation',
                    'Integration endpoints'
                ]
            },
            {
                'name': 'Core Implementation',
                'duration_weeks': 8,
                'key_activities': [
                    'AI model deployment',
                    'Compliance rule configuration',
                    'Data integration',
                    'User interface development'
                ],
                'deliverables': [
                    'Functional AI platform',
                    'Compliance dashboard',
                    'Integrated data flows'
                ]
            },
            {
                'name': 'Testing & Validation',
                'duration_weeks': 3,
                'key_activities': [
                    'System testing',
                    'User acceptance testing',
                    'Performance optimization',
                    'Security validation'
                ],
                'deliverables': [
                    'Test reports',
                    'Performance benchmarks',
                    'Security certification'
                ]
            },
            {
                'name': 'Go-Live & Support',
                'duration_weeks': 2,
                'key_activities': [
                    'Production deployment',
                    'User training',
                    'Knowledge transfer',
                    'Support handover'
                ],
                'deliverables': [
                    'Live system',
                    'Trained users',
                    'Support documentation'
                ]
            }
        ]
        
        return phases
    
    async def _generate_architecture_diagram(
        self, 
        tech_requirements: TechnicalRequirements
    ) -> str:
        """Generate architecture diagram as SVG"""
        
        # This would generate an actual architecture diagram
        # For now, return a placeholder SVG
        svg_diagram = f"""
        <svg width="800" height="600" xmlns="http://www.w3.org/2000/svg">
            <rect width="800" height="600" fill="#0B1020"/>
            <text x="400" y="50" text-anchor="middle" fill="#E6EAF2" font-size="24" font-weight="bold">
                {tech_requirements.scenario} - {tech_requirements.primary_stack} Architecture
            </text>
            
            <!-- User Layer -->
            <rect x="50" y="100" width="700" height="80" fill="#131A2A" stroke="#57D0FF" stroke-width="2" rx="10"/>
            <text x="400" y="145" text-anchor="middle" fill="#E6EAF2" font-size="16">User Interface Layer</text>
            
            <!-- API Gateway -->
            <rect x="200" y="220" width="400" height="60" fill="#1A2332" stroke="#57D0FF" stroke-width="2" rx="10"/>
            <text x="400" y="255" text-anchor="middle" fill="#E6EAF2" font-size="14">API Gateway & Security</text>
            
            <!-- Application Services -->
            <rect x="50" y="320" width="200" height="120" fill="#2A3441" stroke="#57D0FF" stroke-width="1" rx="8"/>
            <text x="150" y="345" text-anchor="middle" fill="#E6EAF2" font-size="12">Compliance Engine</text>
            
            <rect x="300" y="320" width="200" height="120" fill="#2A3441" stroke="#57D0FF" stroke-width="1" rx="8"/>
            <text x="400" y="345" text-anchor="middle" fill="#E6EAF2" font-size="12">AI/ML Services</text>
            
            <rect x="550" y="320" width="200" height="120" fill="#2A3441" stroke="#57D0FF" stroke-width="1" rx="8"/>
            <text x="650" y="345" text-anchor="middle" fill="#E6EAF2" font-size="12">Integration Hub</text>
            
            <!-- Data Layer -->
            <rect x="100" y="480" width="150" height="80" fill="#3A4551" stroke="#57D0FF" stroke-width="1" rx="8"/>
            <text x="175" y="525" text-anchor="middle" fill="#E6EAF2" font-size="12">Database</text>
            
            <rect x="300" y="480" width="150" height="80" fill="#3A4551" stroke="#57D0FF" stroke-width="1" rx="8"/>
            <text x="375" y="525" text-anchor="middle" fill="#E6EAF2" font-size="12">Document Store</text>
            
            <rect x="500" y="480" width="150" height="80" fill="#3A4551" stroke="#57D0FF" stroke-width="1" rx="8"/>
            <text x="575" y="525" text-anchor="middle" fill="#E6EAF2" font-size="12">Cache Layer</text>
            
            <!-- Connections -->
            <line x1="400" y1="180" x2="400" y2="220" stroke="#57D0FF" stroke-width="2"/>
            <line x1="400" y1="280" x2="150" y2="320" stroke="#57D0FF" stroke-width="2"/>
            <line x1="400" y1="280" x2="400" y2="320" stroke="#57D0FF" stroke-width="2"/>
            <line x1="400" y1="280" x2="650" y2="320" stroke="#57D0FF" stroke-width="2"/>
        </svg>
        """
        
        return svg_diagram
    
    def _calculate_competitive_scores(
        self, 
        competitive_data: CompetitiveAnalysis
    ) -> Dict[str, float]:
        """Calculate weighted competitive scores"""
        
        final_scores = {}
        
        for competitor in competitive_data.competitors:
            total_score = 0
            total_weight = 0
            
            for dimension in competitive_data.dimensions:
                score = competitive_data.scores.get(competitor, {}).get(dimension, 0)
                weight = competitive_data.weights.get(dimension, 1)
                
                total_score += score * weight
                total_weight += weight
            
            final_scores[competitor] = total_score / total_weight if total_weight > 0 else 0
        
        return final_scores
    
    async def _generate_competitive_charts(
        self, 
        competitive_data: CompetitiveAnalysis
    ) -> Dict[str, str]:
        """Generate competitive analysis charts"""
        
        charts = {}
        
        # Radar chart
        fig_radar = go.Figure()
        
        for competitor in competitive_data.competitors:
            scores = [
                competitive_data.scores.get(competitor, {}).get(dim, 0)
                for dim in competitive_data.dimensions
            ]
            
            fig_radar.add_trace(go.Scatterpolar(
                r=scores,
                theta=competitive_data.dimensions,
                fill='toself',
                name=competitor
            ))
        
        fig_radar.update_layout(
            polar=dict(
                radialaxis=dict(
                    visible=True,
                    range=[0, 10]
                )
            ),
            showlegend=True,
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)',
            font=dict(color='#E6EAF2')
        )
        
        charts['radar'] = fig_radar.to_html(include_plotlyjs='cdn')
        
        # Bar chart for final scores
        final_scores = self._calculate_competitive_scores(competitive_data)
        
        fig_bar = px.bar(
            x=list(final_scores.keys()),
            y=list(final_scores.values()),
            title="Overall Competitive Scores"
        )
        
        fig_bar.update_layout(
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)',
            font=dict(color='#E6EAF2')
        )
        
        charts['bar'] = fig_bar.to_html(include_plotlyjs='cdn')
        
        return charts


class DocumentExporter:
    """Export proposals to multiple formats"""
    
    def __init__(self):
        self.temp_dir = Path("temp")
        self.temp_dir.mkdir(exist_ok=True)
    
    async def export_to_pdf(self, html_content: str, output_path: str) -> str:
        """Export HTML to PDF using WeasyPrint"""
        
        # Custom CSS for PDF
        pdf_css = CSS(string="""
            @page {
                size: A4;
                margin: 2cm;
                @top-center {
                    content: "DoganAI Compliance - Confidential";
                    font-size: 10px;
                    color: #666;
                }
                @bottom-right {
                    content: "Page " counter(page) " of " counter(pages);
                    font-size: 10px;
                    color: #666;
                }
            }
            
            body {
                font-family: 'Arial', sans-serif;
                line-height: 1.6;
                color: #333;
            }
            
            .page-break {
                page-break-before: always;
            }
            
            .no-break {
                page-break-inside: avoid;
            }
            
            h1, h2, h3 {
                color: #1e3a8a;
                page-break-after: avoid;
            }
            
            table {
                width: 100%;
                border-collapse: collapse;
                margin: 1rem 0;
            }
            
            th, td {
                border: 1px solid #ddd;
                padding: 8px;
                text-align: left;
            }
            
            th {
                background-color: #f2f2f2;
                font-weight: bold;
            }
        """)
        
        html_doc = HTML(string=html_content)
        html_doc.write_pdf(output_path, stylesheets=[pdf_css])
        
        return output_path
    
    async def export_to_pptx(
        self, 
        proposal_data: Dict[str, Any], 
        output_path: str
    ) -> str:
        """Export proposal to PowerPoint presentation"""
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Proposal: {proposal_data.get('title', 'DoganAI Compliance Solution')}"
        subtitle.text = f"Prepared for: {proposal_data.get('client_name', 'Client')}\nDate: {datetime.now().strftime('%B %d, %Y')}"
        
        # Executive Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Executive Summary'
        
        tf = body_shape.text_frame
        tf.text = 'Key Benefits:'
        
        for benefit in proposal_data.get('benefits', []):
            p = tf.add_paragraph()
            p.text = f"• {benefit.get('title', '')}: {benefit.get('description', '')}"
            p.level = 1
        
        # Technical Architecture slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = 'Technical Architecture'
        
        # Add architecture diagram if available
        if 'architecture_diagram' in proposal_data:
            # This would add the actual diagram
            pass
        
        # Implementation Timeline slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        
        title_shape.text = 'Implementation Timeline'
        
        tf = body_shape.text_frame
        tf.text = 'Implementation Phases:'
        
        for phase in proposal_data.get('implementation_phases', []):
            p = tf.add_paragraph()
            p.text = f"• {phase.get('name', '')}: {phase.get('duration_weeks', 0)} weeks"
            p.level = 1
        
        # Save presentation
        prs.save(output_path)
        
        return output_path
    
    async def create_proposal_package(
        self,
        proposal_data: Dict[str, Any],
        formats: List[ExportFormat],
        output_dir: str
    ) -> Dict[str, str]:
        """Create complete proposal package in multiple formats"""
        
        output_files = {}
        base_name = f"proposal_{proposal_data.get('id', 'draft')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        for format_type in formats:
            if format_type == ExportFormat.PDF:
                output_path = f"{output_dir}/{base_name}.pdf"
                await self.export_to_pdf(
                    proposal_data.get('html_content', ''), 
                    output_path
                )
                output_files['pdf'] = output_path
                
            elif format_type == ExportFormat.PPTX:
                output_path = f"{output_dir}/{base_name}.pptx"
                await self.export_to_pptx(proposal_data, output_path)
                output_files['pptx'] = output_path
                
            elif format_type == ExportFormat.JSON:
                output_path = f"{output_dir}/{base_name}.json"
                async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                    await f.write(json.dumps(proposal_data, indent=2, ensure_ascii=False))
                output_files['json'] = output_path
        
        return output_files


class MobileProposalInterface:
    """Mobile-optimized interface for proposal management"""
    
    def __init__(self):
        self.generator = ProposalGenerator(ProposalTemplate())
        self.exporter = DocumentExporter()
    
    async def create_quick_proposal(
        self,
        client_name: str,
        sector: str,
        scenario: str,
        primary_stack: str,
        language: str = 'en'
    ) -> Dict[str, Any]:
        """Create a quick proposal from minimal input"""
        
        # Create simplified client context
        client_context = ClientContext(
            name=client_name,
            sector=sector,
            size="Enterprise",  # Default
            region="KSA",
            current_stack=["Legacy Systems"],
            pain_points=["Manual processes", "Compliance challenges"],
            compliance_requirements=["NCA", "SAMA"] if sector.lower() == 'banking' else ["NCA"],
            budget_range="$500K - $2M",
            timeline="6-12 months",
            decision_makers=[{"name": "CTO", "role": "Technical Decision Maker"}],
            competitors=["Microsoft", "Google", "AWS"],
            kpis=[{"name": "Process Efficiency", "target": "50% improvement"}]
        )
        
        # Create technical requirements
        tech_requirements = TechnicalRequirements(
            scenario=scenario,
            primary_stack=primary_stack,
            secondary_stacks=[],
            deployment_model="hybrid",
            integration_points=["ERP", "CRM", "Legacy Systems"],
            security_requirements=["End-to-end encryption", "Multi-factor authentication"],
            performance_requirements={"response_time": "< 2 seconds", "uptime": "99.9%"},
            scalability_requirements={"users": "1000+", "transactions": "10K/hour"}
        )
        
        # Generate proposal sections
        executive_summary = await self.generator.generate_executive_summary(
            client_context, 
            tech_requirements, 
            language
        )
        
        technical_proposal = await self.generator.generate_technical_proposal(
            client_context,
            tech_requirements,
            language
        )
        
        # Create competitive analysis
        competitive_data = CompetitiveAnalysis(
            dimensions=["Technology", "Support", "Cost", "Implementation Speed"],
            competitors=[primary_stack, "Microsoft", "Google"],
            scores={
                primary_stack: {"Technology": 9, "Support": 8, "Cost": 7, "Implementation Speed": 8},
                "Microsoft": {"Technology": 8, "Support": 7, "Cost": 6, "Implementation Speed": 7},
                "Google": {"Technology": 8, "Support": 6, "Cost": 8, "Implementation Speed": 6}
            },
            weights={"Technology": 0.3, "Support": 0.2, "Cost": 0.3, "Implementation Speed": 0.2},
            rationale={}
        )
        
        competitive_analysis = await self.generator.generate_competitive_analysis(
            competitive_data,
            client_context,
            language
        )
        
        # Combine all sections
        full_proposal = f"""
        <html>
        <head>
            <meta charset="utf-8">
            <title>Proposal for {client_name}</title>
            <style>
                body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                .header {{ text-align: center; margin-bottom: 40px; }}
                .section {{ margin: 30px 0; }}
                .page-break {{ page-break-before: always; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>DoganAI Compliance Solution</h1>
                <h2>Proposal for {client_name}</h2>
                <p>Generated on {datetime.now().strftime('%B %d, %Y')}</p>
            </div>
            
            <div class="section">
                <h2>Executive Summary</h2>
                {executive_summary}
            </div>
            
            <div class="section page-break">
                <h2>Technical Proposal</h2>
                {technical_proposal}
            </div>
            
            <div class="section page-break">
                <h2>Competitive Analysis</h2>
                {competitive_analysis}
            </div>
        </body>
        </html>
        """
        
        proposal_data = {
            'id': str(uuid.uuid4()),
            'client_name': client_name,
            'sector': sector,
            'scenario': scenario,
            'primary_stack': primary_stack,
            'language': language,
            'created_at': datetime.now().isoformat(),
            'status': ProposalStatus.DRAFT,
            'html_content': full_proposal,
            'benefits': self.generator._generate_key_benefits(client_context, tech_requirements),
            'implementation_phases': self.generator._generate_implementation_phases(tech_requirements),
            'roi_projection': self.generator._calculate_roi_projection(client_context)
        }
        
        return proposal_data
    
    async def export_mobile_proposal(
        self,
        proposal_data: Dict[str, Any],
        format_type: ExportFormat = ExportFormat.PDF
    ) -> str:
        """Export proposal optimized for mobile sharing"""
        
        output_dir = "exports"
        Path(output_dir).mkdir(exist_ok=True)
        
        if format_type == ExportFormat.PDF:
            output_path = f"{output_dir}/mobile_proposal_{proposal_data['id']}.pdf"
            return await self.exporter.export_to_pdf(
                proposal_data['html_content'],
                output_path
            )
        elif format_type == ExportFormat.PPTX:
            output_path = f"{output_dir}/mobile_proposal_{proposal_data['id']}.pptx"
            return await self.exporter.export_to_pptx(proposal_data, output_path)
        else:
            # Return base64 encoded JSON for easy mobile transmission
            json_data = json.dumps(proposal_data, ensure_ascii=False)
            return base64.b64encode(json_data.encode()).decode()


# Integration with existing mobile UI
def integrate_proposal_builder_with_mobile():
    """Integration function for adding proposal builder to mobile UI"""
    
    mobile_interface = MobileProposalInterface()
    
    # Add to mobile UI navigation
    proposal_menu = {
        "title": "Proposal Builder",
        "icon": "??",
        "description": "Generate compliance proposals",
        "action": "proposal_builder",
        "features": [
            "Quick proposal generation",
            "Multi-format export (PDF, PPTX)",
            "Competitive analysis",
            "ROI calculations",
            "Arabic/English support"
        ]
    }
    
    return mobile_interface, proposal_menu